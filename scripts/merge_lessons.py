#!/usr/bin/env python3
"""Merge War Horse Lesson 12-15 decks into a single PPTX in order.

Copies shapes, images, backgrounds, and notes from each source deck.
"""

import copy
from pptx import Presentation

INPUT_FILES = [
    "output/WH_Lesson12.pptx",
    "output/WH_Lesson13.pptx",
    "output/WH_Lesson14.pptx",
    "output/WH_Lesson15.pptx",
]
OUTPUT = "output/WH_Lessons12-15_Combined.pptx"


def copy_slide(src_prs, src_slide, dst_prs):
    """Copy a slide from src into dst, preserving images and notes."""
    dst_layout = dst_prs.slide_layouts[0]
    dst_slide = dst_prs.slides.add_slide(dst_layout)

    # Replace shape tree
    src_spTree = src_slide.shapes._spTree
    dst_spTree = dst_slide.shapes._spTree
    dst_spTree.getparent().replace(dst_spTree, copy.deepcopy(src_spTree))

    # Copy background
    src_bg = src_slide.background._element
    dst_bg = dst_slide.background._element
    dst_bg.getparent().replace(dst_bg, copy.deepcopy(src_bg))

    # Copy image parts — map old rIds to new ones and fix references
    rid_map = {}
    for rId, rel in list(src_slide.part.rels.items()):
        if "image" in rel.reltype:
            # relate_to returns the new rId
            new_rId = dst_slide.part.relate_to(rel.target_part, rel.reltype)
            rid_map[rId] = new_rId

    # Update rIds in the copied shape tree if they changed
    if rid_map:
        nsmap = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main",
                 "r": "http://schemas.openxmlformats.org/officeDocument/2006/relationships"}
        for old_rId, new_rId in rid_map.items():
            if old_rId != new_rId:
                for elem in dst_slide.shapes._spTree.iter():
                    for attr_name in ("embed", "{%s}embed" % nsmap["r"],
                                      "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed"):
                        val = elem.get(attr_name)
                        if val == old_rId:
                            elem.set(attr_name, new_rId)

    # Copy notes
    if src_slide.has_notes_slide:
        notes_text = src_slide.notes_slide.notes_text_frame.text
        if notes_text.strip():
            dst_slide.notes_slide.notes_text_frame.text = notes_text

    return dst_slide


def merge():
    base_prs = Presentation(INPUT_FILES[0])

    for fpath in INPUT_FILES[1:]:
        src_prs = Presentation(fpath)
        for slide in src_prs.slides:
            copy_slide(src_prs, slide, base_prs)

    base_prs.save(OUTPUT)
    total = len(base_prs.slides)
    print(f"Merged {len(INPUT_FILES)} decks -> {OUTPUT} ({total} slides)")


if __name__ == "__main__":
    merge()
